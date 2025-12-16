from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)
shapes = slide.shapes

# Colors
green_dark = RGBColor(0x2F, 0x6F, 0x5E)
text_gray = RGBColor(0x2E, 0x2E, 0x2E)

# Title
title_tf = shapes.add_textbox(Inches(1.0), Inches(0.4), Inches(11.333), Inches(1.0)).text_frame
p = title_tf.paragraphs[0]
p.text = "Décisions d'Indexation et Performances"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = green_dark
p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

# Subtitle
sub_tf = shapes.add_textbox(Inches(1.5), Inches(1.4), Inches(10.333), Inches(0.6)).text_frame
p = sub_tf.paragraphs[0]
p.text = "Indexer pour accélérer les requêtes fréquentes et filtrer par proximité/temps/clé."
p.font.size = Pt(18)
p.font.color.rgb = text_gray
p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

# Separator (green rounded rectangle as line)
sep = shapes.add_shape(1, Inches(1.6), Inches(2.2), Inches(11.333), Inches(0.12))
sep.fill.solid()
sep.fill.fore_color.rgb = green_dark
sep.line.color.rgb = green_dark

# Bullets box
bx = shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(7.5), Inches(4.0))
tf = bx.text_frame
tf.word_wrap = True

bullets = [
    ("2dsphere : bin.location — recherches géospatiales (proximité des bacs)."),
    ("Index simple : incident.bacId, vehicle.id, employee.id — accès direct par identifiants."),
    ("Index composé : (routeId, active) — retrouver rapidement RouteActive par itinéraire et état."),
    ("Index temporel / TTL : logs.createdAt (TTL) — purge automatique des logs anciens."),
    ("Index texte : incident.description — recherches textuelles sur descriptions d'incident."),
    ("Index unique : user.username, user.email — garantir unicité et recherches rapides."),
    ("Observabilité : surveiller explain() et tailles d'index ; éviter champs à très faible sélectivité."),
    ("Implémentation (Spring Data / MongoDB) : @Indexed, @CompoundIndex, @GeoSpatialIndexed, @Indexed(expireAfterSeconds=...).")
]

for i, b in enumerate(bullets):
    if i == 0:
        p = tf.paragraphs[0]
        p.text = b
    else:
        p = tf.add_paragraph()
        p.text = b
    p.level = 0
    p.font.size = Pt(16)
    p.font.color.rgb = text_gray

# Footer
f_tf = shapes.add_textbox(Inches(1.0), Inches(6.6), Inches(11.333), Inches(0.6)).text_frame
p = f_tf.paragraphs[0]
p.text = "Réduit la latence des queries critiques — équilibre coût d'écriture vs lectures fréquentes."
p.font.size = Pt(12)
p.font.color.rgb = text_gray
p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

# Save
out = "indexes_slide.pptx"
prs.save(out)
print(f"Saved {out}")
